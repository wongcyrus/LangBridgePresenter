#!/usr/bin/env python3
"""
Speaker Note Generator
Enhances PowerPoint presentations by generating speaker notes using a Supervisor-Tool Multi-Agent System.

Refactored to use Python-based Agent definitions.
"""

import argparse
import asyncio
import logging
import os
import sys
from typing import Dict, Any

import pymupdf  # fitz
from PIL import Image
from pptx import Presentation

from google.adk.runners import InMemoryRunner
from google.genai import types
from google.adk.agents import LlmAgent

# Import Agents
# Ensure the path includes the current directory to find 'agents'
sys.path.append(os.path.dirname(os.path.abspath(__file__)))
from agents.supervisor import supervisor_agent
from agents.analyst import analyst_agent

# Setup logging
logging.basicConfig(
    level=logging.INFO,
    format='%(levelname)s:%(name)s:%(message)s',
    stream=sys.stdout
)
logger = logging.getLogger(__name__)

# Global registry for images
IMAGE_REGISTRY: Dict[str, Image.Image] = {}

async def run_stateless_agent(agent: LlmAgent, prompt: str, image: Image.Image = None) -> str:
    """Helper to run a stateless single-turn agent."""
    runner = InMemoryRunner(agent=agent, app_name=agent.name)
    user_id = "system_user"
    
    parts = [types.Part.from_text(text=prompt)]
    if image:
        parts.append(types.Part.from_image(image=image))

    content = types.Content(role='user', parts=parts)
    
    # Create new session for statelessness
    session = await runner.session_service.create_session(
        app_name=agent.name,
        user_id=user_id,
    )
    
    response_text = ""
    try:
        # Run agent
        for event in runner.run(
            user_id=user_id,
            session_id=session.session_id,
            new_message=content,
        ):
            if getattr(event, "content", None) and event.content.parts:
                part = event.content.parts[0]
                text = getattr(part, "text", "") or ""
                response_text += text
    except Exception as e:
        logger.error(f"Error running agent {agent.name}: {e}")
        return f"Error: {e}"
        
    return response_text.strip()

# Tool wrapper for Analyst
async def call_analyst(image_id: str) -> str:
    """Tool: Analyzes the slide image."""
    logger.info(f"[Tool] call_analyst invoked for image_id: {image_id}")
    image = IMAGE_REGISTRY.get(image_id)
    if not image:
        return "Error: Image not found."
    
    prompt_text = "Analyze this slide image."
    return await run_stateless_agent(analyst_agent, prompt_text, image=image)


async def process_presentation(pptx_path: str, pdf_path: str, course_id: str = None):
    
    logger.info(f"Processing PPTX: {pptx_path}")
    
    # Load files
    prs = Presentation(pptx_path)
    pdf_doc = pymupdf.open(pdf_path)
    limit = min(len(prs.slides), len(pdf_doc))
    
    # Configure Supervisor Tools
    # We append the function tool 'call_analyst' to the existing list of AgentTools
    # supervisor_agent.tools is a list
    if call_analyst not in supervisor_agent.tools:
        supervisor_agent.tools.append(call_analyst)

    # Initialize Supervisor Runner
    supervisor_runner = InMemoryRunner(
        agent=supervisor_agent, 
        app_name="supervisor"
    )

    # Global Context
    presentation_theme = "General Presentation" 
    if course_id:
        try:
            # Dynamically import to avoid circular imports
            project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            if project_root not in sys.path:
                 sys.path.append(project_root)
            from presentation_preloader.utils import course_utils
            course_config = course_utils.get_course_config(course_id)
            if course_config:
                presentation_theme = course_config.get("description") or course_config.get("name") or f"Course {course_id}"
                logger.info(f"Set Presentation Theme from Course: {presentation_theme}")
        except Exception as e:
            logger.warning(f"Failed to fetch course config for {course_id}: {e}")

    previous_slide_summary = "Start of presentation."

    user_id = "supervisor_user"
    session_id = "supervisor_session" 
    
    # Create Supervisor Session
    await supervisor_runner.session_service.create_session(
        app_name="supervisor",
        user_id=user_id,
        session_id=session_id
    )

    for i in range(limit):
        slide_idx = i + 1
        slide = prs.slides[i]
        pdf_page = pdf_doc[i]
        
        logger.info(f"--- Processing Slide {slide_idx} ---")
        
        # 1. Setup Slide Context
        existing_notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
             existing_notes = slide.notes_slide.notes_text_frame.text.strip()

        # Register Image
        image_id = f"slide_{slide_idx}"
        pix = pdf_page.get_pixmap(dpi=150)
        IMAGE_REGISTRY[image_id] = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

        # 2. Prompt Supervisor
        supervisor_prompt = (
            f"Here is Slide {slide_idx}.\n"
            f"Existing Notes: \"{existing_notes}\"\n"
            f"Image ID: \"{image_id}\"\n"
            f"Previous Slide Summary: \"{previous_slide_summary}\"\n"
            f"Theme: \"{presentation_theme}\"\n\n"
            f"Please proceed with the workflow."
        )

        content = types.Content(
            role='user',
            parts=[types.Part.from_text(text=supervisor_prompt)]
        )

        # 3. Run Supervisor Loop
        final_response = ""
        
        try:
            for event in supervisor_runner.run(
                user_id=user_id,
                session_id=session_id,
                new_message=content,
            ):
                if getattr(event, "content", None) and event.content.parts:
                     part = event.content.parts[0]
                     text = getattr(part, "text", "") or ""
                     final_response += text
        except Exception as e:
            logger.error(f"Error in supervisor loop: {e}")

        final_response = final_response.strip()
        logger.info(f"Final Note for Slide {slide_idx}: {final_response[:50]}...")
        
        # 4. Update PPTX
        if not slide.has_notes_slide:
             try:
                 slide.notes_slide # This might create it in some versions or fail
             except:
                 pass # Handling logic depends on pptx version

        try:
            slide.notes_slide.notes_text_frame.text = final_response
        except Exception as e:
            logger.error(f"Could not write note: {e}")
            
        # 5. Update Context
        previous_slide_summary = final_response[:200]

        # Cleanup Image
        del IMAGE_REGISTRY[image_id]

    # Save
    output_path = pptx_path.replace(".pptx", "_enhanced.pptx")
    prs.save(output_path)
    logger.info(f"Saved enhanced deck to: {output_path}")

def main():
    parser = argparse.ArgumentParser(description="Generate Speaker Notes with Supervisor Agent")
    parser.add_argument("--pptx", required=True, help="Path to input PPTX")
    parser.add_argument("--pdf", required=True, help="Path to input PDF")
    parser.add_argument("--course-id", help="Optional: Course ID to fetch theme context")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.pptx) or not os.path.exists(args.pdf):
        print("Error: Input files not found.")
        return

    asyncio.run(process_presentation(args.pptx, args.pdf, args.course_id))

if __name__ == "__main__":
    main()
