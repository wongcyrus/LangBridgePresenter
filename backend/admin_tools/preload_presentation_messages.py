#!/usr/bin/env python3
"""
Pre-generate presentation_messages from a PPTX slide and cache them
for fast retrieval by generate_presentation_message().

This tool writes to TWO locations:
1. Cache entries in `xiaoice_presentation_cache` collection for each language
2. Config document `xiaoice_config/messages` with presentation_messages map

Usage:
  python preload_presentation_messages.py \
    --pptx /path/to/deck.pptx \
    --languages en,zh \
    --template "Welcome to {title}" \
    --slide 1 \
    --context "quarterly review"

Notes:
  - Cache entries use key "v1:{language}:{hash}" format matching firestore_utils.py
  - Hash is first 12 chars of SHA256 of normalized context
  - Empty context uses "default" instead of hash
  - No translation is performed; same template used for all languages
"""

import argparse
import hashlib
import os
from typing import Dict, List

from google.cloud import firestore
from pptx import Presentation


def parse_languages(s: str) -> List[str]:
    return [x.strip() for x in s.split(",") if x.strip()]


def _normalize_context(context: str) -> str:
    """Normalize context by trimming and collapsing whitespace."""
    if not context:
        return ""
    # Collapse all whitespace runs to a single space and strip ends
    return " ".join(str(context).split())


def _cache_key(language_code: str, context: str) -> str:
    """Build a stable, short cache key safe for Firestore doc IDs.

    Uses lowercase language + 12-char SHA256 of normalized context.
    Avoids very long document IDs and ensures consistent lookups.
    Must match the logic in firestore_utils.py
    """
    lang = (language_code or "").strip().lower() or "unknown"
    norm_ctx = _normalize_context(context)
    if not norm_ctx:
        return f"v1:{lang}:default"
    digest = hashlib.sha256(norm_ctx.encode("utf-8")).hexdigest()[:12]
    return f"v1:{lang}:{digest}"


def get_slide_text(prs: Presentation, slide_index: int) -> str:
    if slide_index < 1 or slide_index > len(prs.slides):
        raise IndexError(
            f"slide index {slide_index} out of range (1..{len(prs.slides)})"
        )
    slide = prs.slides[slide_index - 1]
    # Prefer first text-containing shape as title/summary
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = (shape.text or "").strip()
            if text:
                return text
    return ""


def build_message(title: str, template: str) -> str:
    return template.replace("{title}", title or "").strip()


def cache_message(
    db: firestore.Client, language_code: str, context: str, message: str
):
    """Write a cache entry for generate_presentation_message lookup."""
    cache_key = _cache_key(language_code, context)
    norm_ctx = _normalize_context(context)
    cache_ref = db.collection("xiaoice_presentation_cache").document(cache_key)
    cache_ref.set({
        "message": message,
        "language_code": (language_code or "").strip().lower(),
        "context": norm_ctx,
        "context_hash": cache_key.rsplit(":", 1)[-1],
        "updated_at": firestore.SERVER_TIMESTAMP
    })


def update_config(db: firestore.Client, messages: Dict[str, str]):
    """Update presentation_messages in config document."""
    doc_ref = db.collection("xiaoice_config").document("messages")
    doc = doc_ref.get()
    current = doc.to_dict() if doc.exists else {}
    merged = dict(current or {})
    merged["presentation_messages"] = messages
    merged["updated_at"] = firestore.SERVER_TIMESTAMP
    doc_ref.set(merged)


def main():
    parser = argparse.ArgumentParser(
        description="Preload presentation message cache from PPTX"
    )
    parser.add_argument("--pptx", required=True, help="Path to PPTX file")
    parser.add_argument(
        "--languages", default="en", help="Comma list, e.g. en,zh"
    )
    parser.add_argument(
        "--template",
        default="Welcome to {title}",
        help="Message template; supports {title}",
    )
    parser.add_argument(
        "--slide",
        type=int,
        default=1,
        help="Slide index to extract title/text from (1-based)",
    )
    parser.add_argument(
        "--context",
        default="",
        help="Context string to include in cache key (optional)",
    )

    args = parser.parse_args()

    if not os.path.exists(args.pptx):
        raise FileNotFoundError(args.pptx)

    languages = parse_languages(args.languages)
    prs = Presentation(args.pptx)
    title = get_slide_text(prs, args.slide)
    message = build_message(title, args.template)

    db = firestore.Client(database="xiaoice")

    # Write cache entries for each language
    for lang in languages:
        cache_message(db, lang, args.context, message)
        cache_key = _cache_key(lang, args.context)
        print(f"Cached '{cache_key}': {message}")

    # Update config document
    messages = {lang: message for lang in languages}
    update_config(db, messages)
    print(
        f"\nUpdated presentation_messages in config "
        f"for {len(languages)} languages"
    )


if __name__ == "__main__":
    main()
